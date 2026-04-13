"""Build the LLM-fuzzing review deck as a .pptx from the collected A/B data.

Sources drawn from (already on disk):
  - docs/AB_RE2_REPORT.md   (headline numbers, per-file tables, bug list)
  - docs/STATUS.md          (pipeline shape + current state)

Output: docs/slides/llm_fuzzing_review.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = REPO_ROOT / "docs" / "slides" / "llm_fuzzing_review.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_text(tf, text: str, size: int = 18, bold: bool = False, align=None):
    tf.clear()
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold


def _add_title(slide, title: str, subtitle: str | None = None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    _set_text(tb.text_frame, title, size=32, bold=True)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.5))
        _set_text(sb.text_frame, subtitle, size=16)


def _add_body(slide, lines: list[str], top: float = 1.8, size: int = 18,
              height: float = 5.3, left: float = 0.6):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(12.2), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        p.space_after = Pt(6)


def _add_table(slide, rows: list[list[str]], top: float = 2.0, left: float = 0.8,
               width: float = 11.5, height: float = 4.0, header_bold: bool = True,
               font_size: int = 16):
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    table = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                   Inches(width), Inches(height)).table
    for r_i, row in enumerate(rows):
        for c_i in range(n_cols):
            cell = table.cell(r_i, c_i)
            val = row[c_i] if c_i < len(row) else ""
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(font_size)
            if r_i == 0 and header_bold:
                run.font.bold = True
    return table


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # blank

    # --- Slide 1: Title ---
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.8))
    _set_text(tb.text_frame,
              "LLM-Guided Fuzzing Seed Corpora",
              size=44, bold=True, align=PP_ALIGN.CENTER)
    sb = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.3), Inches(1.2))
    _set_text(sb.text_frame,
              "First A/B experiment on RE2 +\n"
              "the regex-format prompt rewrite that made the result trustworthy",
              size=22, align=PP_ALIGN.CENTER)
    meta = s.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.0))
    _set_text(meta.text_frame,
              "Shreyas Ganesh   |   2026-04-13\n"
              "Target: RE2  |  Model: llama-3.1-8b-instruct (UF LiteLLM proxy)  |  Spend: ~$0.01",
              size=14, align=PP_ALIGN.CENTER)

    # --- Slide 2: Research question ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "The research question",
               "Given a stable pipeline, does gap-targeted synthesis beat source-only?")
    _add_body(s, [
        "Exp1  (gap-targeted):   prompt = harness + coverage gaps + a few upstream tests",
        "Exp2  (source-only):    prompt = harness + library source code, no tests, no gaps",
        "",
        "Shared: 3-test RE2 fixture (Regexp.BigRef, NamedCaptures, CaptureNames),",
        "llama-3.1-8b-instruct, 3 samples per cell, T=0.7, top_p=0.95, identical context otherwise.",
        "",
        "Evaluation: raw seed-corpus edge/line coverage on a standalone seed_replay binary",
        "(same LLVMFuzzerTestOneInput as the libFuzzer target, but without injected main()).",
        "",
        "This deck is about seed-corpus quality, not 24h libFuzzer-campaign coverage.",
        "That number needs cluster time we don't yet have.",
    ], size=18)

    # --- Slide 3: Pipeline shape ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Pipeline shape (exp1 and exp2 share everything but the prompt)")
    _add_body(s, [
        "1.  Phase 1   —  build coverage-instrumented target, extract upstream tests,",
        "                 compute per-test + union baseline coverage, emit coverage_gaps.json.",
        "",
        "2.  Phase 2   —  LLM predicts hard branches from the test context (exp1 only).",
        "",
        "3.  Phase 3   —  LLM synthesises inputs (3 samples, T=0.7). Samples cache-salted",
        "                 so the 3 calls don't collide on prompt_hash.",
        "",
        "4.  Evaluation —  seed_replay + llvm-profdata + llvm-cov-15 export → CoverageProfile",
        "                  JSON → analysis.scripts.ab_coverage_diff → line/edge set-diff +",
        "                  Jaccard + per-file tables.",
        "",
        "exp1 uses gap-targeted templates; exp2 uses source-only templates with a hard",
        "assert_no_tests guard on the rendered prompt.",
    ], size=17)

    # --- Slide 4: Stabilization recap ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Recap — the 2026-04-12 stabilization pass",
               "Detail: docs/WORK_SO_FAR.md §2-4")
    _add_body(s, [
        "Before:  every synthesis call hit max_tokens, multi-sample runs silently cache-collided,",
        "         schema placeholders were copied verbatim into degenerate loops.",
        "",
        "Landed:",
        "  •  cache key includes max_tokens + cache_salt (3-sample diversity is real now)",
        "  •  streaming loop-detector (two signals: window dominance + unique-window ratio)",
        "  •  5 prompt templates rewritten with real example values + OUTPUT RULES",
        "  •  strict=False cache read, extra='ignore' on PredictionResult",
        "  •  29 new unit tests, full suite 106+ green",
        "",
        "Result: 4-cell sanity run completes in ~3 min for ~$0.02.",
        "That set the stage for the first real A/B — which surfaced its own bug class.",
    ], size=17)

    # --- Slide 5: First A/B — bugs enabling it ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "First A/B surfaced five more bugs")
    _add_body(s, [
        "1.  Fixture tests pointed at re2_test.cc, but our build only compiled regexp_test.",
        "    Fix: swap fixture to tests that exist in the built binary.",
        "",
        "2.  RE2's util/test.h ignores --gtest_filter; all registered tests always run.",
        "    Not a bug — a constraint. Treat fixture as a combined-test baseline.",
        "",
        "3.  CoverageProfile JSON round-trip failed — writer used field names, reader",
        "    wanted the true/false aliases. Fix: populate_by_name=True on FrozenConfig.",
        "",
        "4.  llvm-cov-15 export --format=json rejected (JSON is the default).",
        "    Fix: drop the flag in run_test_coverage.py + measure_coverage.py.",
        "",
        "5.  Loop-aborted samples produced 0 seeds even when the first array entries were",
        "    complete. Fix: _salvage_objects() brace-matching scan for any balanced {…}",
        "    sub-object carrying content_b64.",
    ], size=16)

    # --- Slide 6: Bytes-format A/B result ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "First A/B result — bytes-format prompt (2026-04-12)",
               "Prompts asked for base64-encoded bytes; parser decoded straight to seed files.")
    _add_table(s, [
        ["metric", "exp1 (gap)", "exp2 (source)", "union", "intersection"],
        ["seeds produced", "5", "10", "—", "—"],
        ["loop-aborted samples", "3 / 3", "2 / 3", "—", "—"],
        ["edges covered", "475", "513", "530", "458"],
        ["lines covered", "1243", "1314", "1337", "1220"],
        ["edges ONLY in this cell", "17", "55", "—", "—"],
        ["Jaccard (edges)", "—", "—", "0.864", "—"],
    ], top=2.0, height=3.2, font_size=16)
    _add_body(s, [
        "Headline at the time:   \"exp2 wins, 513 > 475.\"",
        "Caveat at the time:     seed-count asymmetry (5 vs 10 after salvage)",
        "                        + 5 of 6 synthesis samples loop-aborted.",
    ], top=5.4, size=16)

    # --- Slide 7: Diagnosis ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Why bytes-format was a bad fit",
               "RE2 harness layout vs what the prompt was asking for")
    _add_body(s, [
        "RE2 libFuzzer harness (simplified):",
        "",
        "    if (size < 3 || size > 64) return 0;",
        "    uint8_t flag_bytes[2] = { data[0], data[1] };",
        "    std::string pattern(data + 2, size - 2);",
        "    RE2::Compile(pattern, flags_from(flag_bytes));",
        "",
        "The harness expects  [2 flag bytes] [UTF-8 regex string].",
        "We asked the LLM for arbitrary base64 bytes and hoped they'd look like a regex.",
        "",
        "For a frontier model: maybe. For llama-3.1-8b: the model emits 100× copies",
        "of the same 48-byte placeholder until the loop detector fires.",
        "",
        "The \"loop-rate\" was largely prompt/target mismatch, not model capability.",
    ], size=16)

    # --- Slide 8: Regex-format rewrite ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Fix — rewrite prompts to emit regex text (2026-04-13)")
    _add_body(s, [
        "New Jinja templates:",
        "  •  prediction/prompts/input_synthesis_regex.j2        (gap-targeted)",
        "  •  synthesis/prompts/source_only_synthesis_regex.j2  (source-only)",
        "",
        "New parser  parse_regex_response()  in synthesis/scripts/parse_synthesis.py:",
        "  •  extracts  { \"regexes\": [ { \"regex\": str, \"target_gaps\": [...], … }, … ] }",
        "  •  prepends 2 sha256-seeded flag bytes keyed by (target, sample, idx, regex)",
        "  •  UTF-8 encode with codepoint-boundary-safe truncation",
        "  •  clip total bytes to RE2's [3, 64] window, dedupe",
        "",
        "New CLI flag  --input-format {bytes,regex}  on both generators routes",
        "the template + parser pair. Old bytes outputs archived under *_bytes_v1/.",
    ], size=16)

    # --- Slide 9: Regex-format A/B result ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Second A/B result — regex-format prompt (2026-04-13)",
               "Same target, same fixture, same model. Only prompt + parser changed.")
    _add_table(s, [
        ["metric", "exp1 (gap)", "exp2 (source)", "union", "intersection"],
        ["seeds produced", "20", "30", "—", "—"],
        ["ok samples / total", "2 / 3", "3 / 3", "—", "—"],
        ["edges covered", "1243", "1133", "1308", "1068"],
        ["lines covered", "2530", "2385", "2661", "2254"],
        ["edges ONLY in this cell", "175", "65", "—", "—"],
        ["Jaccard (edges)", "—", "—", "0.817", "—"],
    ], top=2.0, height=3.2, font_size=16)
    _add_body(s, [
        "Coverage per cell roughly tripled.    Loop-abort rate 5/6 → 1/6.",
        "Exp1 (gap-targeted) now wins by +110 edges and 175 vs 65 exclusive.",
    ], top=5.4, size=17)

    # --- Slide 10: Bytes vs Regex side-by-side ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Bytes vs regex — what the rewrite actually changed")
    _add_table(s, [
        ["metric", "bytes (04-12)", "regex (04-13)", "change"],
        ["exp1 seeds", "5", "20", "+4×"],
        ["exp2 seeds", "10", "30", "+3×"],
        ["exp1 edges", "475", "1243", "+2.6×"],
        ["exp2 edges", "513", "1133", "+2.2×"],
        ["exp1 exclusive edges", "17", "175", "+10×"],
        ["exp2 exclusive edges", "55", "65", "+1.2×"],
        ["loop-aborted samples", "5 / 6", "1 / 6", "−80%"],
        ["headline winner", "exp2 (+38)", "exp1 (+110)", "flipped"],
    ], top=1.9, height=4.8, font_size=16)

    # --- Slide 11: Where exp1 wins uniquely ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Where exp1 wins uniquely in the regex run",
               "175 edges reached by exp1 that exp2 misses — top files:")
    _add_table(s, [
        ["file", "unique edges (exp1 only)"],
        ["re2/parse.cc", "86"],
        ["re2/regexp.cc", "50"],
        ["re2/simplify.cc", "14"],
        ["re2/compile.cc", "9"],
        ["re2/prog.cc", "6"],
        ["re2/re2.cc", "3"],
        ["re2/dfa.cc", "3"],
        ["util/logging.h", "2"],
    ], top=2.0, height=4.0, font_size=17)
    _add_body(s, [
        "Parser / simplifier / compiler path. Exactly the branches the gap-targeted",
        "prompt highlights. exp2 sees parser source but lacks the coverage annotation.",
    ], top=6.1, size=16)

    # --- Slide 12: The takeaway ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "The takeaway")
    _add_body(s, [
        "1.  Prompt format matters as much as prompt content.",
        "    A content-correct prompt in the wrong format looks identical to",
        "    \"the model can't do the task.\" We almost concluded the bytes-format",
        "    A/B favored exp2 when the real story was the harness/format mismatch.",
        "",
        "2.  Match the prompt to the target's native input idiom.",
        "    Regex for RE2,  SQL for sqlite,  binary header for libpng.",
        "    This is not a one-size template the way the research plan implied.",
        "",
        "3.  At matched format, the gap-targeted hypothesis holds on RE2 at this scale.",
        "    But: n=3 samples, 1 target, 1 model — not yet a statistical claim.",
    ], size=18)

    # --- Slide 12b: Generalization follow-up — question ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Generalization follow-up — is exp1's +110 transferable?",
               "The 2026-04-13 A/B measures exp1 on the same fixture its gap list came from.")
    _add_body(s, [
        "Hypothesis (user):  exp2 (source-only) should generalize better because",
        "                    exp1's gap list is a flashlight on a known neighborhood,",
        "                    not a transferable skill.",
        "",
        "Three probes, all llama-8b, all regex format, all under $1 LLM spend:",
        "",
        "  P0   random baseline          — does either LLM cell beat chance?",
        "  A    held-out source subset   — does exp1's win transfer to files its",
        "                                  gap list never pointed at?",
        "  B    7-cell prompt ablation   — which piece of exp1's context is",
        "                                  actually doing the work?",
        "",
        "Experiment C (1h × 3-trial libFuzzer campaigns) deferred — needs build/fuzzer/",
        "rebuild + ~9 CPU-hours; A+B evidence was judged sufficient.",
    ], size=17)

    # --- Slide 12c: P0 random baseline ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "P0 — Random baseline establishes the floor",
               "30 seeds of [2 random flag bytes][random ASCII 1-62], seeded random.Random(42)")
    _add_table(s, [
        ["cell", "seeds", "edges", "Δ vs random"],
        ["exp1 (gap-targeted)", "20", "1243", "+263  (+27%)"],
        ["exp2 (source-only)",  "30", "1133", "+153  (+16%)"],
        ["random baseline",     "30",  "980",  "—"],
    ], top=2.0, height=2.2, font_size=18)
    _add_body(s, [
        "Both LLM cells clear the floor meaningfully — the \"+110 edges\" claim",
        "is not a rounding artefact on top of noise.",
        "",
        "But: exp1's +110 lead over exp2 is only ≈42% of exp1's +263 lead over random.",
        "Most of the above-random signal is reachable by exp2's recipe already.",
    ], top=4.6, size=17)

    # --- Slide 12d: Experiment A — held-out source subset ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Experiment A — held-out source-file subset",
               "Split RE2 into A (parser, seen by exp1 gaps) vs B (execution, held out).")
    _add_body(s, [
        "Set A (gap list restricted here):   re2/parse.cc, re2/regexp.cc,",
        "                                     re2/simplify.cc, re2/tostring.cc",
        "Set B (held out, coverage measured): re2/{compile,prog,dfa,nfa,onepass,",
        "                                     bitstate,re2}.cc + util/{rune,strutil}.cc",
    ], top=1.5, size=15, height=1.4)
    _add_table(s, [
        ["cell", "seeds", "B-edges", "Δ vs exp2 on B"],
        ["exp1_full",       "20", "599", "+3"],
        ["exp2_source",     "30", "596", "0"],
        ["exp1_heldout",    "30", "581", "−15"],
        ["random baseline", "30", "557", "−39"],
    ], top=3.2, height=2.1, font_size=17)
    _add_body(s, [
        "+110 in-distribution  →  +3 on held-out files  →  −15 when exp1's gaps",
        "are restricted to a disjoint file set. exp2's recipe transfers; exp1's does not.",
    ], top=5.6, size=17)

    # --- Slide 12e: Experiment B — prompt ablation ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Experiment B — 7-cell prompt ablation",
               "Which piece of exp1's extra context (gaps vs tests) is actually load-bearing?")
    _add_table(s, [
        ["cell", "gaps", "tests", "source", "edges", "Δ vs exp1_full"],
        ["exp2_plus_gaps",   "✓", "",  "✓", "1250", "+7"],
        ["exp1_full",        "✓", "✓", "",  "1243", "0"],
        ["exp2_plus_tests",  "",  "✓", "✓", "1210", "−33"],
        ["exp2_source",      "",  "",  "✓", "1133", "−110"],
        ["exp1_tests_only",  "",  "✓", "",  "1093", "−150"],
        ["random",           "—", "—", "—",  "980", "−263"],
        ["exp1_gaps_only*",  "✓", "",  "",   "879", "−364"],
    ], top=1.9, height=3.8, font_size=14)
    _add_body(s, [
        "1. Source is the load-bearing context (all source cells ≥ 1133).",
        "2. Gaps amplify source:  exp2_plus_gaps beats exp2_source by +117  (largest single-variable effect).",
        "3. Tests alone ≈ source alone  (1093 vs 1133).",
        "*exp1_gaps_only: 2/3 samples loop-aborted on llama — dense gaps w/o source/tests derails the model.",
    ], top=5.9, size=13)

    # --- Slide 12f: Verdict ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Verdict — hypothesis supported, with a sharper recipe")
    _add_body(s, [
        "The user's hypothesis — \"exp2 generalizes better than exp1\" — is",
        "supported with nuance:",
        "",
        "  • On the same files the gap list points at, exp1 wins decisively  (+110).",
        "  • On held-out files within the same target, exp1's advantage collapses  (+3).",
        "  • When exp1's gaps are restricted to a disjoint file set, exp1 loses  (−15).",
        "",
        "The most efficient single recipe on this fixture is exp2_plus_gaps",
        "(source code + coverage gaps, no tests) — +7 over exp1_full, +117 over exp2_source.",
        "",
        "Clean take: the right baseline is not exp1 or exp2 — it's",
        "source + cheap coverage annotations. exp1's +110 was mostly the annotation doing",
        "work exp2's source already covered, and on held-out files the annotation stops",
        "generalizing.",
    ], size=16)

    # --- Slide 13: What this does NOT prove ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "What this does NOT prove")
    _add_body(s, [
        "•  Single target, single model.",
        "   Held-out subset test is intra-target; cross-target transfer is still open.",
        "   Need a second target in its own native idiom before claiming the recipe generalizes.",
        "",
        "•  No libFuzzer campaign numbers.",
        "   This deck measures the seed corpus directly; the research-relevant",
        "   metric is 24h-campaign coverage starting from each seed set (Experiment C, deferred).",
        "",
        "•  Not yet exercised against a frontier model.",
        "   The exp1_gaps_only collapse is llama-specific — may disappear on GPT-4o / Claude.",
        "   A frontier rerun on just that cell (~$1) would make or break the \"gaps are",
        "   load-bearing\" claim.",
        "",
        "•  n≈20–30 per cell supports direction, not a p-value.",
        "   Rigorous stats would need n≥30 per cell and multiple random A/B file splits.",
        "",
        "•  Fine-tuning path (Phase 4) is dry-run-only until A100 access.",
    ], size=16)

    # --- Slide 14: Reproducibility ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Reproducibility — the 2026-04-13 regex A/B",
               "Full bash in docs/AB_RE2_REPORT.md § Reproducibility")
    _add_body(s, [
        "# exp1 (gap-targeted)",
        "UTCF_LITELLM_URL=https://api.ai.it.ufl.edu UTCF_LLM_RPM=12 \\",
        ".venv/bin/python -m synthesis.scripts.generate_inputs \\",
        "    --target re2 --model llama-3.1-8b-instruct --samples 3 --experiment exp1 \\",
        "    --input-format regex \\",
        "    --dataset-root dataset/fixtures/re2_ab \\",
        "    --results-root dataset/fixtures/re2_ab/phase3_results",
        "",
        "# exp2 (source-only)",
        "UTCF_LITELLM_URL=https://api.ai.it.ufl.edu UTCF_LLM_RPM=12 \\",
        ".venv/bin/python -m synthesis.scripts.generate_source_inputs \\",
        "    --target re2 --model llama-3.1-8b-instruct --samples 3 \\",
        "    --source-max-files 4 --source-token-budget 14000 --max-tokens 8192 \\",
        "    --input-format regex \\",
        "    --results-root dataset/fixtures/re2_ab/exp2_results",
        "",
        "# coverage + diff — see AB report for measure_coverage + ab_coverage_diff",
        "",
        "Total LLM spend: ~$0.007.    Wall-clock: under 2 minutes.",
    ], size=14)

    # --- Slide 15: Updated decision asks ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Decision asks")
    _add_body(s, [
        "1.  Authorize a frontier-model rerun of Experiment B.",
        "    GPT-4o / Claude Sonnet on the 7-cell ablation. ~$10.",
        "    Specifically tests whether exp1_gaps_only's collapse is model-specific",
        "    and whether exp2_plus_gaps stays ahead of exp1_full.",
        "",
        "2.  Commit to per-target prompt templates for any new target in the A/B set.",
        "    RE2 shows this is non-optional  (regex for RE2, SQL for sqlite, …).",
        "",
        "3.  Promote exp2_plus_gaps to the reference recipe going forward.",
        "    Source + coverage gaps (no tests) is the efficient frontier on this fixture.",
        "",
        "4.  Unblock Experiment C: build/fuzzer/ rebuild + ~9 CPU-hours for 1h×3-trial",
        "    campaigns. Or the earlier 29,440 CPU-hour 24h-campaign ask — either way the",
        "    regex-format seeds are the right input to that run.",
    ], size=18)

    # --- Slide 16: Backing material ---
    s = prs.slides.add_slide(blank)
    _add_title(s, "Backing material")
    _add_body(s, [
        "docs/STATUS.md",
        "    Living handoff doc. Read first for current state, active work,",
        "    planned next steps, and the 1-screen pipeline overview.",
        "",
        "docs/AB_RE2_REPORT.md",
        "    Full A/B writeup — headline numbers, bytes-run archive, reproducibility bash.",
        "",
        "docs/claude_code_plan_v3.md  /  docs/research_document_v3.md",
        "    Authoritative execution plan + research design.",
        "",
        "analysis/scripts/ab_coverage_diff.py",
        "    Turns two CoverageProfile JSONs into set-wise line/edge diff + Jaccard",
        "    + per-file tables. Source of every number in slide 10.",
        "",
        "dataset/fixtures/re2_ab/ab_coverage/",
        "    The exp1.json, exp2.json, ab_coverage_diff.{json,md} this deck is built on.",
    ], size=16)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}")


if __name__ == "__main__":
    build()
